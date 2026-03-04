"""
Generate PowerPoint presentation from PRESENTATION_MATERIAL.md.
Requires: pip install python-pptx
Run: python generate_presentation_pptx.py
Output: GNN3_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import os

OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "GNN3_Presentation.pptx")


def add_title_slide(prs, title, subtitle=""):
    layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_bullet_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
    return slide


def add_table_slide(prs, title, headers, rows):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.6)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    # Table
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_width = Inches(9)
    col_width = table_width / n_cols
    table = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.1), table_width, Inches(0.4 * n_rows)).table
    for i, h in enumerate(headers):
        table.cell(0, i).text = str(h)
        table.cell(0, i).text_frame.paragraphs[0].font.bold = True
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            if c_idx < n_cols:
                table.cell(r_idx + 1, c_idx).text = str(val)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(prs, "GNN-Based Power Flow Surrogate", "IEEE 34-Bus Distribution Feeder")

    # 2. Problem
    add_bullet_slide(prs, "Problem & Motivation", [
        "Need fast voltage prediction for real-time grid applications",
        "OpenDSS power-flow solves are iterative and slow",
        "Goal: Replace with a GNN surrogate for instant inference",
    ])

    # 3. Approach
    add_bullet_slide(prs, "Approach", [
        "Message-passing GNN on phase-level graph (95 nodes, 184 edges)",
        "OpenDSS for ground-truth simulation and dataset generation",
        "PyTorch Geometric for GNN training",
        "5-min resolution (288 timesteps per day)",
    ])

    # 4. Datasets
    add_table_slide(prs, "Datasets (5 Types)", 
        ["#", "Name", "Features", "Target", "Samples"],
        [
            ["1", "Original", "p_load, q_load, p_pv (3)", "vmag_pu", "~192k"],
            ["2", "Injection", "p_inj, q_inj (2)", "vmag_pu", "~192k"],
            ["3", "Load-type", "13 (M1/M2/M4/M5, dist, balance)", "vmag_pu", "~192k"],
            ["4", "Delta-V", "14 (Load-type + vmag_zero)", "vmag_delta_pu", "57,600"],
            ["5", "Delta-V 5×", "14 (same)", "vmag_delta_pu", "57,600"],
        ])

    # 5. Dataset Generation
    add_bullet_slide(prs, "Dataset Generation", [
        "Scenario sampling: P_load, Q_load, P_PV from BASELINE × RANGES",
        "24h profiles: load shape mL[t], irradiance mPV[t] from CSV",
        "Time sampling: select_times_three_profiles (load, PV, net bins)",
        "Delta-V fix: Use mL[t], mPV[t] in training to match test-time structure",
    ])

    # 6. Architecture
    add_bullet_slide(prs, "GNN Architecture (PFIdentityGNN)", [
        "Node embeddings: learnable per-node (identity)",
        "Edge features: R, X from OpenDSS",
        "Message passing: EdgeIdentityMP (MLP on h_j + edge_attr + edge_emb)",
        "Readout: MLP → per-node voltage prediction",
        "Hyperparameters: node_emb, edge_emb, h_dim, num_layers, phase_onehot",
    ])

    # 7. Exploration
    add_table_slide(prs, "Architecture Exploration (7 Blocks)",
        ["Block", "Configs", "Datasets", "Best"],
        [
            ["1", "20 nominees", "Original, Injection, Load-type", "light_xwide"],
            ["2", "60 nominees", "Same 3", "—"],
            ["3", "10 promising", "Same 3", "—"],
            ["4", "10 to beat Block 3", "Same 3", "—"],
            ["5", "8 from Block 4", "Same 3", "light_emb_h96_phase_onehot_depth3_h112"],
            ["6", "9 best Load-type", "Delta-V", "light_xwide_emb_phase_onehot"],
            ["7", "9 best", "Delta-V 5×", "light_xwide_emb_phase_onehot"],
        ])

    # 8. Final Models
    add_table_slide(prs, "Final 7 Models (Best per Dataset)",
        ["Block", "Dataset", "Architecture", "Target"],
        [
            ["1", "Original", "medium (h=32, 4L)", "vmag_pu"],
            ["2", "Injection", "deep (h=64, 4L)", "vmag_pu"],
            ["3", "Load-type", "light_emb_h96 (h=96, 2L)", "vmag_pu"],
            ["4", "Load-type", "light_emb_h96_phase_onehot_depth3", "vmag_pu"],
            ["5", "Load-type", "light_emb_h96_phase_onehot_depth3_h112", "vmag_pu"],
            ["6", "Delta-V", "light_xwide_emb_phase_onehot", "vmag_delta_pu"],
            ["7", "Delta-V 5×", "light_xwide_emb_phase_onehot", "vmag_delta_pu"],
        ])

    # 9. Results
    add_bullet_slide(prs, "Results", [
        "Best full voltage (V): Block 5 — MAE ≈ 0.006 pu",
        "Best delta-V: Block 6/7 — MAE ≈ 0.003–0.005 pu",
        "24h overlay plots: OpenDSS vs GNN at observed node",
        "PV sensitivity: 0.7×, 1.0×, 1.3× penetration",
    ])

    # 10. Delta-V Retrain & Voltage Profile Plotting
    add_bullet_slide(prs, "Delta-V Retrain & Voltage Profile Plotting", [
        "run_gnn3_deltav_only_train: retrain blocks 6 & 7 on new 24h-profile datasets",
        "plot_voltage_profile_by_operation_point: multi-node, multi-operation-point 24h plots",
        "Explores oscillation-free operation points (P_BASE, Q_BASE, PV_BASE)",
    ])

    # 11. Error Analysis
    add_bullet_slide(prs, "Error Analysis: Spatial and Temporal Breakdown", [
        "Spatial: Per-bus error to identify high-error locations",
        "Temporal: Per-timestep error to locate spikes",
        "Used to understand where the model fails",
    ])

    # 12. Key Details
    add_bullet_slide(prs, "Key Details", [
        "Exploration: 30% data; final training: 100% data",
        "Phase onehot: added in Blocks 4–7 for phase-aware prediction",
        "Batched GNN inference: 288 timesteps in one forward pass",
        "Timing: OpenDSS profile vs full GNN pipeline (zero-PV + build_x + forward)",
    ])

    # 13. Summary
    add_bullet_slide(prs, "Summary", [
        "5 datasets: Original, Injection, Load-type, Delta-V, Delta-V 5×",
        "7 exploration blocks, 100+ architecture configs",
        "7 best models, one per dataset type",
        "Delta-V train–test alignment fix (24h profile)",
        "PV sensitivity, error analysis, voltage profile plotting",
    ])

    # 14. Thank you
    add_title_slide(prs, "Thank You", "Questions?")

    prs.save(OUTPUT_PATH)
    print(f"[SAVED] {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
