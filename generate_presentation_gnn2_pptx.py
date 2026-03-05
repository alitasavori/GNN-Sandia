"""
Generate PowerPoint presentation for GNN2 notebook.
Requires: pip install python-pptx
Run: python generate_presentation_gnn2_pptx.py
Output: GNN2_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import os

OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "GNN2_Presentation.pptx")


def add_title_slide(prs, title, subtitle=""):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_bullet_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
    return slide


def add_table_slide(prs, title, headers, rows):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    left, top, width, height = Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_width = Inches(9)
    table = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.1), table_width, Inches(0.4 * n_rows)).table
    for i, h in enumerate(headers):
        table.cell(0, i).text = str(h)
        table.cell(0, i).text_frame.paragraphs[0].font.bold = True
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            if c_idx < n_cols:
                table.cell(r_idx + 1, c_idx).text = str(val)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(prs, "GNN2: Dataset Generation & Initial Training", "IEEE 34-Bus Power Flow Surrogate")

    # 2. Overview
    add_bullet_slide(prs, "GNN2 Notebook Overview", [
        "Colab setup: clone repo, install torch, torch-geometric, opendssdirect",
        "Dataset generation: 5 datasets from OpenDSS power-flow solves",
        "Graph sanity checks: bidirectionality, edges, node features",
        "Initial GNN training: Original, Injection, Load-type datasets",
        "Unified config exploration: light, medium, heavy, deep per dataset",
    ])

    # 3. Dataset Generation Pipeline
    add_bullet_slide(prs, "Dataset Generation Pipeline (Step by Step)", [
        "1. Read OpenDSS model, set daily mode (5-min, 288 pts/day)",
        "2. Find load/irradiance profile CSVs from DSS",
        "3. Build bus-phase node list (95 nodes)",
        "4. Extract static phase edges (R, X, C) → 184 edges",
        "5. Sample scenarios: P_load, Q_load, P_PV from BASELINE × RANGES",
        "6. Select timestamps: 3 profiles (load, PV, net) with anchors + equal-pop bins",
        "7. Apply snapshot: device setpoints + noise, solve OpenDSS",
        "8. Store node features and voltage targets",
    ])

    # 4. Five Datasets
    add_table_slide(prs, "Five Datasets",
        ["#", "Name", "Output", "Features", "Target"],
        [
            ["1", "Original", "datasets_gnn2/original", "p_load, q_load, p_pv (3)", "vmag_pu"],
            ["2", "Injection", "datasets_gnn2/injection", "p_inj, q_inj (2)", "vmag_pu"],
            ["3", "Load-type", "datasets_gnn2/loadtype", "13 (M1/M2/M4/M5, dist, balance)", "vmag_pu"],
            ["4", "Delta-V", "datasets_gnn2/deltav", "14 (+ vmag_zero)", "vmag_delta_pu"],
        ])

    # 5. Output Files
    add_bullet_slide(prs, "Output Files (per dataset)", [
        "gnn_edges_phase_static.csv — edges with R_full, X_full, C_full",
        "gnn_node_index_master.csv — node names and indices",
        "gnn_sample_meta.csv — sample_id, scenario_id, P_load, Q_load, P_pv, m_load, m_pv",
        "gnn_node_features_and_targets.csv — features + vmag_pu (or vmag_delta_pu)",
    ])

    # 6. Time Selection
    add_bullet_slide(prs, "Time Selection: 3 Profiles + Anchors", [
        "Three profiles: load (mL), PV (mPV), net (P_load×mL − P_pv×mPV)",
        "Split K snapshots across profiles (~1/3 each)",
        "Anchors: min and max time index per profile (extreme operating points)",
        "Equal-population bins: rank-based, same count per bin",
        "Ensures coverage of load-high, PV-high, net-stress regions",
    ])

    # 7. Graph Sanity Check
    add_bullet_slide(prs, "Graph Sanity Check", [
        "Bidirectionality: every edge (a→b) has reverse (b→a)",
        "Transformer/regulator edges: 832↔888, 814↔814r, 852↔852r",
        "One-sample node features: verify feature columns",
        "Nominal power balance: sum_load − sum_pv ≈ p_sys_balance",
        "Delta-V summary: vmag_zero_pv_pu, vmag_delta_pu (dataset 4)",
    ])

    # 8. GNN Architecture (PFIdentityGNN)
    add_bullet_slide(prs, "GNN Architecture (PFIdentityGNN)", [
        "Message-passing GNN on phase-level graph",
        "Node embeddings: learnable per-node (identity)",
        "Edge features: R_full, X_full from OpenDSS",
        "Message: MLP(h_j + edge_attr + edge_emb)",
        "Readout: MLP → per-node voltage (vmag_pu or vmag_delta_pu)",
    ])

    # 9. Nominal vs Realized Power Comparison
    add_bullet_slide(prs, "Nominal vs Realized Power Comparison", [
        "Runs OpenDSS for 10 samples; compares nominal (set) vs realized (power-flow) values",
        "Loads: per-load and per-node; PV: nominal vs realized; Capacitors; Grid (TotalPower)",
        "Sum check: verifies realized P/Q sums ≈ 0 (power balance)",
        "Fix: %Cutin=0 on PV systems (default 2% cut-in caused zero output at low irradiance)",
    ])

    # 10. Initial Training
    add_table_slide(prs, "Initial Training (per dataset)",
        ["Dataset", "Features", "Checkpoint"],
        [
            ["Original", "3 (p_load, q_load, p_pv)", "pf_identity_gnn_best.pt"],
            ["Injection", "2 (p_inj, q_inj)", "pf_identity_gnn_inj_best.pt"],
            ["Load-type", "13 (load types, dist, balance)", "pf_identity_gnn_loadtype_best.pt"],
        ])

    # 11. VMAG-only & Embedding Inspection
    add_bullet_slide(prs, "VMAG-only Training & Embedding Inspection", [
        "VMAG-only training: outputs vmag_pu only (angle removed); smaller, faster model",
        "Embedding inspection: print node/edge identity embeddings (z_i, r_e) for PFIdentityGNN",
        "Verifies learned per-node and per-edge representations",
    ])

    # 12. Unified Config Exploration
    add_bullet_slide(prs, "Unified Config Exploration", [
        "4 configs: light, medium, heavy, deep (varying h_dim, num_layers)",
        "Run each on all 3 datasets (Original, Injection, Load-type)",
        "12 runs total — find best dataset × architecture combination",
        "Load-type typically wins; feeds into GNN3 exploration",
    ])

    # 13. Summary
    add_bullet_slide(prs, "GNN2 Summary", [
        "5 datasets generated from OpenDSS with 24h profiles",
        "Nominal vs realized comparison; %Cutin=0 fix for PV",
        "Graph sanity checks ensure valid phase-level topology",
        "PFIdentityGNN trained on Original, Injection, Load-type",
        "Unified config exploration (light, medium, heavy, deep) → feeds GNN3",
    ])

    # 14. Thank you
    add_title_slide(prs, "Thank You", "Questions?")

    prs.save(OUTPUT_PATH)
    print(f"[SAVED] {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
